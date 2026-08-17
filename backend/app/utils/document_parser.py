from io import BytesIO

from fastapi import UploadFile
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract


ALLOWED_EXTENSIONS = {
    ".pdf",
    ".txt",
    ".docx",
    ".pptx",
    ".png",
    ".jpg",
    ".jpeg",
}


TESSERACT_PATH = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

pytesseract.pytesseract.tesseract_cmd = TESSERACT_PATH


def get_file_extension(filename: str) -> str:
    return "." + filename.rsplit(".", 1)[-1].lower()


async def extract_text(file: UploadFile) -> str:
    if not file.filename:
        raise ValueError("File name is required")

    extension = get_file_extension(file.filename)

    if extension not in ALLOWED_EXTENSIONS:
        raise ValueError(
            "Unsupported file type. "
            "Supported formats: PDF, TXT, DOCX, PPTX, PNG, JPG, JPEG"
        )

    file_content = await file.read()

    if not file_content:
        raise ValueError("Uploaded file is empty")

    # --------------------------------------------------
    # TXT
    # --------------------------------------------------

    if extension == ".txt":
        try:
            text = file_content.decode("utf-8")
        except UnicodeDecodeError:
            text = file_content.decode("utf-8", errors="ignore")

        if not text.strip():
            raise ValueError("Text file contains no readable content")

        return text.strip()

    # --------------------------------------------------
    # PDF
    # --------------------------------------------------

    if extension == ".pdf":
        try:
            reader = PdfReader(BytesIO(file_content))
        except Exception as exc:
            raise ValueError(
                f"Could not read PDF file: {exc}"
            )

        pages = []

        for page in reader.pages:
            text = page.extract_text()

            if text:
                pages.append(text.strip())

        extracted_text = "\n\n".join(pages)

        if not extracted_text.strip():
            raise ValueError(
                "Could not extract text from PDF. "
                "The PDF may contain scanned images."
            )

        return extracted_text.strip()

    # --------------------------------------------------
    # DOCX
    # --------------------------------------------------

    if extension == ".docx":
        try:
            document = Document(BytesIO(file_content))
        except Exception as exc:
            raise ValueError(
                f"Could not read DOCX file: {exc}"
            )

        paragraphs = []

        for paragraph in document.paragraphs:
            text = paragraph.text.strip()

            if text:
                paragraphs.append(text)

        # Also extract text from tables
        for table in document.tables:
            for row in table.rows:
                row_text = []

                for cell in row.cells:
                    cell_text = cell.text.strip()

                    if cell_text:
                        row_text.append(cell_text)

                if row_text:
                    paragraphs.append(" | ".join(row_text))

        extracted_text = "\n\n".join(paragraphs)

        if not extracted_text.strip():
            raise ValueError(
                "Could not extract readable text from DOCX file"
            )

        return extracted_text.strip()

    # --------------------------------------------------
    # PPTX
    # --------------------------------------------------

    if extension == ".pptx":
        try:
            presentation = Presentation(BytesIO(file_content))
        except Exception as exc:
            raise ValueError(
                f"Could not read PPTX file: {exc}"
            )

        slides_text = []

        for slide_number, slide in enumerate(
            presentation.slides,
            start=1,
        ):
            slide_parts = []

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()

                    if text:
                        slide_parts.append(text)

            if slide_parts:
                slides_text.append(
                    f"Slide {slide_number}\n"
                    + "\n".join(slide_parts)
                )

        extracted_text = "\n\n".join(slides_text)

        if not extracted_text.strip():
            raise ValueError(
                "Could not extract readable text from PPTX file"
            )

        return extracted_text.strip()

    # --------------------------------------------------
    # IMAGE OCR
    # --------------------------------------------------

    if extension in {".png", ".jpg", ".jpeg"}:
        try:
            image = Image.open(BytesIO(file_content))
        except Exception as exc:
            raise ValueError(
                f"Could not read image file: {exc}"
            )

        try:
            extracted_text = pytesseract.image_to_string(
                image,
                lang="eng",
            )
        except Exception as exc:
            raise ValueError(
                f"OCR processing failed: {exc}"
            )

        if not extracted_text.strip():
            raise ValueError(
                "Could not extract readable text from image"
            )

        return extracted_text.strip()

    raise ValueError("Unsupported file type")